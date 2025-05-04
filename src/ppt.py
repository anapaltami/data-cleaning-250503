from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd

# load cleaned data
df = pd.read_csv('output/cleaned_data.csv')

# create presentation
presentation = Presentation()
blank_slide_layout = presentation.slide_layouts[6]  # Blank Slide

# styles for the slides
# Catppuccin Mocha Theme
bg_color = RGBColor(30, 30, 46)         # base
title_color = RGBColor(205, 214, 244)   # text
content_color = RGBColor(166, 173, 200) # subtext0
box_color = RGBColor(49, 50, 68)        # surface1
accent_color = RGBColor(203, 166, 247)  # mauve

def add_background_color(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

def add_title(slide, title):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Fira Code'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = title_color

def add_content(slide, content, top=Inches(1.5)):
    body_box = slide.shapes.add_textbox(Inches(0.7), top, Inches(8), Inches(5))
    tf = body_box.text_frame
    tf.clear()  # remove default empty paragraph

    for line in content.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Fira Code'
        p.font.size = Pt(20)
        p.font.color.rgb = content_color

def add_mono_content(slide, content, top=Inches(1.5), left=Inches(0.7), width=Inches(8), height=Inches(3.5)):
    # add rounded rectangle shape behind the content
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = box_color  # surface1
    bg_shape.line.color.rgb = accent_color  # mauve
    bg_shape.shadow.inherit = False  # remove default shadow if any

    # add text box on top of the rounded rectangle
    body_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = body_box.text_frame
    tf.clear()

    for line in content.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Fira Code'
        p.font.size = Pt(14)
        p.font.color.rgb = content_color  # dark gray

def add_accent_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7),
        Inches(1.2),
        Inches(7.6),
        Inches(0.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()  # no border


# slide 1: Title Slide
slide = presentation.slides.add_slide(blank_slide_layout)
add_background_color(slide)
add_title(slide, "Data Cleaning Project")
add_accent_bar(slide)
add_content(slide, "Prepared by: Anastasia Parks Altamirano\nDate: 2025-05-03")

# slide 2: Objectives
slide = presentation.slides.add_slide(blank_slide_layout)
add_background_color(slide)
add_title(slide, "Objectives")
add_accent_bar(slide)
content = (
    "1. Load the data from 'data/FAKE-datasetPII.xlsx'\n"
    "2. Rename columns for easier handling\n"
    "3. Extract ZIP from ADDRESS before dropping\n"
    "4. Remove PII columns\n"
    "5. Scrub the data for PII\n"
    "6. Save the cleaned data to output folder"
)
add_content(slide, content)

# slide 3: Original Dataset
slide = presentation.slides.add_slide(blank_slide_layout)
add_background_color(slide)
add_title(slide, "Original Dataset Overview")
add_accent_bar(slide)
content = (
    "The original dataset contains various columns, \nincluding PII such as:\n"
    "- Full Name\n"
    "- SSN\n"
    "- Credit Card Number\n"
    "- Address\n"
    "- Date of Birth\n"
    "- Phone Number\n"
    "\n"
    "The goal is to protect user privacy while \npreserving useful information."
)
add_content(slide, content)

# slide 4: Data Cleaning Steps
slide = presentation.slides.add_slide(blank_slide_layout)
add_background_color(slide)
add_title(slide, "Data Cleaning Steps")
add_accent_bar(slide)
content = (
    "1. Load the data from 'data/FAKE-datasetPII.xlsx'\n"
    "2. Rename columns for easier handling\n"
    "3. Extract ZIP from ADDRESS before dropping\n"
    "4. Remove PII columns\n"
    "5. Scrub the data for PII\n"
    "6. Save the cleaned data to 'output/cleaned_data.csv'"
)
add_content(slide, content)

# slide 5: Tools Used
slide = presentation.slides.add_slide(blank_slide_layout)
add_background_color(slide)
add_title(slide, "Tools Used")
add_accent_bar(slide)
content = (
    "1. Python\n"
    "2. Pandas\n"
    "3. Scrubadub\n"
    "4. Python-pptx\n"
    "\n"
    "These tools were used to load, clean, and present \nthe process."
)
add_content(slide, content)

# slide 6: Cleaned Data Preview
slide = presentation.slides.add_slide(blank_slide_layout)
add_background_color(slide)
add_title(slide, "Cleaned Data Preview")
add_accent_bar(slide)
df_preview = df.head().to_string(index=False)
content = (
    "The cleaned data preview:\n\n"
    f"{df_preview}\n\n"
    "The cleaned data has been saved to 'output/cleaned_data.csv'."
)
add_mono_content(slide, content)

# slide 7: Card Type Chart
chart_path = 'output/card_type_chart.png'
slide = presentation.slides.add_slide(blank_slide_layout)
add_background_color(slide)
add_title(slide, "Card Type Distribution")
add_accent_bar(slide)
add_content(slide, "The chart below shows the distribution of card types \nin the dataset:\n")
slide.shapes.add_picture(chart_path, Inches(1.5), Inches(2.8), width=Inches(6.5))

# save presentation
presentation.save('slides/data_cleaning_presentation.pptx')
print("Presentation created and saved as 'slides/data_cleaning_presentation.pptx'")